
# https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
# pip install python-pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

def extract_code(path_to_presentation = 'static/PlayFair Ciphertext Encryption.pptx'):
    prs = Presentation(path_to_presentation)
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
                    
    for i in text_runs:
        print(i, end='\n'+'-'*20+'\n')

    return text_runs

# extract_code()


def set_title(
    title_slide_layout = 0,
    path_to_presentation = 'static/hello world.pptx',
    title_text = "Hello, World!",
    subtitle_text = "python-pptx was here!",
    ):
    
    title_slide_layout = prs.slide_layouts[title_slide_layout]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text
    subtitle.text = subtitle_text
    prs.save(path_to_presentation)

# set_title()


def bullet_level(
    bullet_slide_layout = 1,
    path_to_presentation = 'static/nested bullets.pptx',
    title_shape_text = 'Adding a Bullet Slide',
    tf_text = 'Find the bullet slide layout',
    p_text_l1 = 'Use _TextFrame.text for first bullet',
    # p_text_l2 = 'Use _TextFrame.add_paragraph() for subsequent bullets',
    ):
    
    bullet_slide_layout = prs.slide_layouts[bullet_slide_layout]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = title_shape_text

    tf = body_shape.text_frame
    tf.text = tf_text

    p = tf.add_paragraph()
    p.level = 1
    p.text = p_text_l1

    # p = tf.add_paragraph()
    # p.level = 2
    # p.text = p_text_l2

    prs.save(path_to_presentation)

# bullet_level()


def change_fonts(
    blank_slide_layout = 1,
    path_to_presentation = 'static/fonts.pptx',
    tf_text = "This is text inside a textbox",
    p_text_bold = "This is a second paragraph that's bold",
    p_font_bold = True,
    p_text_size = "This is a third paragraph that's big",
    p_font_size = 40,
    ):
    
    blank_slide_layout = prs.slide_layouts[blank_slide_layout]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)

    tf = txBox.text_frame
    tf.text = tf_text

    p = tf.add_paragraph()
    p.text = p_text_bold
    p.font.bold = p_font_bold

    p = tf.add_paragraph()
    p.text = p_text_size
    p.font.size = Pt(p_font_size)

    prs.save(path_to_presentation)

# change_fonts()


def set_image(
    blank_slide_layout = 1,
    prs_save = 'static/image.pptx',
    img_path = 'static/logo.png',
    ):

    img_path = img_path
    blank_slide_layout = prs.slide_layouts[blank_slide_layout]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = Inches(2)
    # pic = slide.shapes.add_picture(img_path, left, top)

    # left = Inches(5)
    height = Inches(5)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)

    prs.save(prs_save)

# set_image()


def shape_steps(
    title_only_slide_layout = 1,
    prs_save = 'static/shapes.pptx',
    shapes_title_text = 'Adding an AutoShape',
    ):

    title_only_slide_layout = prs.slide_layouts[title_only_slide_layout]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = shapes_title_text

    left = Inches(0.93)  # 0.93" centers this overall set of shapes
    top = Inches(3.0)
    width = Inches(1.75)
    height = Inches(1.0)

    shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
    shape.text = 'Step 1'

    left = left + width - Inches(0.4)
    width = Inches(2.0)  # chevrons need more width for visual balance

    for n in range(2, 6):
        shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
        shape.text = 'Step %d' % n
        left = left + width - Inches(0.4)

    prs.save(prs_save)

# shape_steps()


def set_tables(
    title_only_slide_layout = 1,
    prs_save = 'static/table.pptx',
    shapes_title_text = 'Adding a Table',
    ):

    title_only_slide_layout = prs.slide_layouts[title_only_slide_layout]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = shapes_title_text

    rows = cols = 2
    left = top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(0.8)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(4.0)

    # write column headings
    table.cell(0, 0).text = 'Foo'
    table.cell(0, 1).text = 'Bar'

    # write body cells
    table.cell(1, 0).text = 'Baz'
    table.cell(1, 1).text = 'Qux'

    prs.save(prs_save)

# set_tables()
